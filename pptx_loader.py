from pptx import Presentation
from langchain_core.documents import Document
import os
import logging
from utils import setup_logger

logger = setup_logger("PptxLoader")

def load_pptx(file_path: str) -> list[Document]:
    """
    Extract text and slide content from PPTX files.
    """
    file_name = os.path.basename(file_path)
    documents = []
    
    try:
        prs = Presentation(file_path)
        logger.info(f"Opening PPTX: {file_name}, total slides: {len(prs.slides)}")
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                metadata = {
                    "source": file_name,
                    "file_name": file_name,
                    "page_number": i + 1,
                    "file_type": "pptx",
                    "content_type": "slide_text",
                    "method": "python-pptx"
                }
                documents.append(Document(page_content="\n".join(slide_text), metadata=metadata))
                
    except Exception as e:
        logger.error(f"Error loading PPTX {file_path}: {e}")
        
    return documents
